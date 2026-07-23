"""Build the JIVA overview PowerPoint from captured stills.

Usage:
    python3 -m venv .venv && .venv/bin/pip install python-pptx Pillow
    .venv/bin/python build_deck.py

Reads  ./stills/*.png  (produced by capture-stills.mjs)
Writes ../Data/JIVA_Overview.pptx
"""

import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

DIR = os.path.dirname(os.path.abspath(__file__))
STILLS = os.path.join(DIR, "stills")
JPG = os.path.join(STILLS, "_jpg")
OUT = os.path.join(os.path.dirname(DIR), "Data", "JIVA_Overview.pptx")

GREEN = RGBColor(0x2A, 0x61, 0x30)   # Jiva Green
LIME = RGBColor(0xD5, 0xE2, 0x74)    # Jiva Lime
SAGA = RGBColor(0xDD, 0xEE, 0xDE)    # Whispering Saga
DEW = RGBColor(0xF3, 0xF9, 0xF3)     # Dewdrop Glow
INK = RGBColor(0x1A, 0x21, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEAD_FONT = "Lexend"   # Mustachio stand-in per VideoPlan.md
BODY_FONT = "Inter"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def prep_image(name, width=1920):
    """Downscale a still to a deck-friendly JPEG, return path + aspect ratio."""
    os.makedirs(JPG, exist_ok=True)
    src = os.path.join(STILLS, f"{name}.png")
    dst = os.path.join(JPG, f"{name}.jpg")
    img = Image.open(src).convert("RGB")
    ratio = img.height / img.width
    if img.width > width:
        img = img.resize((width, round(width * ratio)), Image.LANCZOS)
    img.save(dst, "JPEG", quality=85, optimize=True)
    return dst, ratio


def blank_slide(prs, bg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT):
    """runs: list of (text, font_name, size_pt, color, bold) tuples, one per paragraph."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, font, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def add_picture_framed(slide, path, ratio, left, top, width):
    pic = slide.shapes.add_picture(path, left, top, width=width, height=Emu(int(width * ratio)))
    pic.line.color.rgb = SAGA
    pic.line.width = Pt(1.25)
    return pic


def add_card(slide, left, top, width, height, fill=WHITE):
    from pptx.enum.shapes import MSO_SHAPE

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = SAGA
    card.line.width = Pt(1)
    card.shadow.inherit = False
    return card


def screen_slide(prs, name, headline, body, ratio_map):
    slide = blank_slide(prs, WHITE)
    add_text(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.7),
             [(headline, HEAD_FONT, 30, GREEN, True)])
    add_text(slide, Inches(0.7), Inches(1.0), Inches(11.9), Inches(0.6),
             [(body, BODY_FONT, 14, INK, False)])
    path, ratio = ratio_map[name]
    img_w = Inches(9.6)
    img_h = img_w * ratio
    left = Emu(int((SLIDE_W - img_w) / 2))
    add_picture_framed(slide, path, ratio, left, Inches(1.75), img_w)
    return slide


def main():
    names = ["intake", "vitality-map", "biomarker", "action-plan", "vitality-map-2"]
    ratio_map = {}
    for n in names:
        path, ratio = prep_image(n)
        ratio_map[n] = (path, ratio)
    logo = os.path.join(STILLS, "logo.png")
    logo_ratio = Image.open(logo).height / Image.open(logo).width

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title
    slide = blank_slide(prs, DEW)
    logo_w = Inches(3.2)
    slide.shapes.add_picture(logo, Emu(int((SLIDE_W - logo_w) / 2)), Inches(1.7),
                             width=logo_w, height=Emu(int(logo_w * logo_ratio)))
    add_text(slide, Inches(1.5), Inches(4.1), Inches(10.33), Inches(0.9),
             [("Precision Wellness · Personalized to You", HEAD_FONT, 32, GREEN, True)],
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.5), Inches(5.0), Inches(10.33), Inches(0.5),
             [("Product overview", BODY_FONT, 16, INK, False)], align=PP_ALIGN.CENTER)

    # 2. The problem
    slide = blank_slide(prs, WHITE)
    add_text(slide, Inches(1.2), Inches(2.3), Inches(11), Inches(1.4),
             [("A lab report is data, not guidance.", HEAD_FONT, 40, GREEN, True)])
    add_text(slide, Inches(1.2), Inches(3.9), Inches(9.8), Inches(1.6),
             [("A standard blood panel returns more than 100 values. Most people cannot act on a single one of them.",
               BODY_FONT, 18, INK, False),
              ("JIVA turns those numbers into a plan.", BODY_FONT, 18, GREEN, True)])

    # 3. How JIVA works
    slide = blank_slide(prs, DEW)
    add_text(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
             [("How JIVA works", HEAD_FONT, 34, GREEN, True)])
    steps = [
        ("1", "Questionnaire",
         "Lifestyle, nutrition, activity, medical and family history shape everything that follows."),
        ("2", "Panel and blood draw",
         "One draw at a partner lab covers 69 to 100+ biomarkers, chosen from a base panel plus targeted add-ons."),
        ("3", "Clinical engine",
         "Physician-led analysis maps every value to 10 functional body systems and ranks what matters."),
        ("4", "Vitality Map",
         "A personal report with a biological age and concrete food, supplement, exercise and yoga actions."),
    ]
    card_w, card_h, gap = Inches(2.95), Inches(4.1), Inches(0.2)
    total = card_w * 4 + gap * 3
    x0 = Emu(int((SLIDE_W - total) / 2))
    for i, (num, title, body) in enumerate(steps):
        left = Emu(int(x0 + i * (card_w + gap)))
        add_card(slide, left, Inches(1.75), card_w, card_h)
        add_text(slide, Emu(int(left + Inches(0.25))), Inches(2.0), Inches(0.8), Inches(0.7),
                 [(num, HEAD_FONT, 34, GREEN, True)])
        add_text(slide, Emu(int(left + Inches(0.25))), Inches(2.75), Emu(int(card_w - Inches(0.5))), Inches(0.9),
                 [(title, HEAD_FONT, 18, INK, True)])
        add_text(slide, Emu(int(left + Inches(0.25))), Inches(3.6), Emu(int(card_w - Inches(0.5))), Inches(2.0),
                 [(body, BODY_FONT, 12.5, INK, False)])

    # 4. Intake screen
    screen_slide(prs, "intake",
                 "Personalization starts before the blood draw",
                 "A six-part intake covers lifestyle, nutrition, activity, medical history, family history and goals. Answers steer the panel selection and every later recommendation.",
                 ratio_map)

    # 5. Vitality Map screen
    screen_slide(prs, "vitality-map",
                 "One map of your health",
                 "69 lab values become 10 color-coded body systems, a range breakdown, and a biological age of 30.8 against a calendar age of 29.",
                 ratio_map)

    # 6. Biomarker + Action Plan side by side
    slide = blank_slide(prs, WHITE)
    add_text(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.7),
             [("Every number explained. Every action justified.", HEAD_FONT, 30, GREEN, True)])
    add_text(slide, Inches(0.7), Inches(1.0), Inches(11.9), Inches(0.6),
             [("Each biomarker sits on its reference range in plain language, and every food, supplement and exercise ties back to her specific values.",
               BODY_FONT, 14, INK, False)])
    pair_w = Inches(6.0)
    for i, n in enumerate(["biomarker", "action-plan"]):
        path, ratio = ratio_map[n]
        left = Emu(int(Inches(0.55) + i * (pair_w + Inches(0.25))))
        add_picture_framed(slide, path, ratio, left, Inches(1.85), pair_w)

    # 7. Retest screen
    screen_slide(prs, "vitality-map-2",
                 "Six months later, measurably better",
                 "A follow-up draw compares every system against the baseline. Progress is visible, and the plan updates to maintenance where she has improved.",
                 ratio_map)

    # 8. Why JIVA
    slide = blank_slide(prs, DEW)
    add_text(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
             [("Why JIVA", HEAD_FONT, 34, GREEN, True)])
    pillars = [
        ("Built for Latin America",
         "A Spanish-first experience with WhatsApp scheduling, local payment rails and locally available foods in every plan."),
        ("A physician-led engine",
         "Recommendations anchored to your own numbers, not generic wellness advice."),
        ("Nicoya Blue Zone heritage",
         "Longevity science from one of only five places on earth where people demonstrably live longest."),
    ]
    card_w, card_h, gap = Inches(3.9), Inches(3.6), Inches(0.25)
    total = card_w * 3 + gap * 2
    x0 = Emu(int((SLIDE_W - total) / 2))
    for i, (title, body) in enumerate(pillars):
        left = Emu(int(x0 + i * (card_w + gap)))
        add_card(slide, left, Inches(2.0), card_w, card_h)
        add_text(slide, Emu(int(left + Inches(0.3))), Inches(2.35), Emu(int(card_w - Inches(0.6))), Inches(1.0),
                 [(title, HEAD_FONT, 19, GREEN, True)])
        add_text(slide, Emu(int(left + Inches(0.3))), Inches(3.35), Emu(int(card_w - Inches(0.6))), Inches(1.9),
                 [(body, BODY_FONT, 13, INK, False)])

    # 9. See it in action
    slide = blank_slide(prs, GREEN)
    add_text(slide, Inches(1.5), Inches(2.9), Inches(10.33), Inches(1.2),
             [("See it in action", HEAD_FONT, 44, WHITE, True)], align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.5), Inches(4.1), Inches(10.33), Inches(0.6),
             [("Live product demo", BODY_FONT, 18, LIME, False)], align=PP_ALIGN.CENTER)

    prs.save(OUT)
    size_mb = os.path.getsize(OUT) / 1e6
    print(f"wrote {OUT} ({size_mb:.1f} MB, {len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
